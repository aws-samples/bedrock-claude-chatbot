import streamlit as st
import boto3
from botocore.config import Config
import os
import pandas as pd
import time
import json
import io
import re
import openpyxl
from python_calamine import CalamineWorkbook
from openpyxl.cell import Cell
import plotly.io as pio
from openpyxl.worksheet.cell_range import CellRange
from docx.table import _Cell
from pptx import Presentation
from botocore.exceptions import ClientError
from textractor import Textractor
from textractor.data.constants import TextractFeatures
from textractor.data.text_linearization_config import TextLinearizationConfig
import pytesseract
from PIL import Image
import PyPDF2
import chardet
from docx import Document as DocxDocument
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.document import Document
from docx.text.paragraph import Paragraph
from docx.table import Table as DocxTable
import concurrent.futures
from functools import partial
import random
from agent.chat_agent import build_chat_agent, StreamlitCallbackHandler, SubAgentWorkingsPanel
from agent.sessions import (wrap_attached_documents, read_display_history,
                            list_user_sessions, dataset_uris_in_window,
                            generated_uris_in_window)
from agent.web_tools import WEB_TOOLS
from agent.data_analysis import (AthenaSparkSessionManager, CodeSessionManager,
                                 make_data_analysis_tool)
from agent.doc_generator import make_document_generator_tool
from agent.workers import make_worker_agents_tool
from urllib.parse import urlparse
import plotly.graph_objects as go
config = Config(
    read_timeout=600,  # Read timeout parameter
    retries=dict(
        max_attempts=10  # Handle retries
    )
)

st.set_page_config(initial_sidebar_state="auto")

# Read app configurations
with open('config.json','r',encoding='utf-8') as f:
    config_file = json.load(f)
# pricing info
with open('pricing.json','r',encoding='utf-8') as f:
    pricing_file = json.load(f)
# Bedrock Model info
with open('model_id.json','r',encoding='utf-8') as f:
    model_info = json.load(f)

DYNAMODB = boto3.resource('dynamodb')
COGNITO = boto3.client('cognito-idp')
S3 = boto3.client('s3')
DYNAMODB_TABLE = config_file["DynamodbTable"]
BUCKET = config_file["Bucket_Name"]
OUTPUT_TOKEN = config_file["max-output-token"]
S3_DOC_CACHE_PATH = config_file["document-upload-cache-s3-path"]
TEXTRACT_RESULT_CACHE_PATH = config_file["AmazonTextract-result-cache"]
LOAD_DOC_IN_ALL_CHAT_CONVO = config_file["load-doc-in-chat-history"]
CHAT_HISTORY_LENGTH = config_file["chat-history-loaded-length"]
DYNAMODB_USER = config_file["UserId"]
REGION = config_file["region"]
USE_TEXTRACT = config_file["AmazonTextract"]
CSV_SEPERATOR = config_file["csv-delimiter"]
INPUT_BUCKET = config_file["input_bucket"]
INPUT_S3_PATH = config_file["input_s3_path"]
INPUT_EXT = tuple(f".{x}" for x in config_file["input_file_ext"].split(','))
SESSION_STORAGE = config_file.get("session-storage", "local")  # local | s3 | dynamodb | agentcore
# AgentCore Memory resource id (required when session-storage is "agentcore"); users
# bring their own resource - creating one is out of the app's scope.
AGENTCORE_MEMORY_ID = config_file.get("agentcore-memory-id", "")
# Data-analysis sub-agent (agent-as-tool). Runtime "python" -> AgentCore Code
# Interpreter; "pyspark" -> Athena for Apache Spark (workgroup must be PySpark engine v3).
DATA_ANALYSIS_MODEL = config_file.get("data-analysis-model", "sonnet-4.6")
CODE_INTERPRETER_ID = config_file.get("code-interpreter-id", "")
CODE_INTERPRETER_TIMEOUT = config_file.get("code-interpreter-session-timeout", 3600)
ATHENA_WORKGROUP = config_file.get("athena-work-group-name", "")
# Worker-agent swarm ("Research Workers" tool): parallel ephemeral web-research agents
WORKER_AGENT_MODEL = config_file.get("worker-agent-model", "haiku-4.5")
WORKER_AGENT_MAX_WORKERS = config_file.get("worker-agent-max-workers", 6)
# Capability lists are derived from the model registry (model_id.json specs) instead
# of hand-maintained: reasoning dialect / vision / tool support live per-model there.
MODEL_DISPLAY_NAME = list(model_info.keys())
HYBRID_MODELS = [k for k, v in model_info.items() if v.get("reasoning")]
NON_VISION_MODELS = [k for k, v in model_info.items() if not v.get("vision", True)]
NON_TOOL_SUPPORTING_MODELS = [k for k, v in model_info.items() if not v.get("tools", True)]

# Orchestrator system prompts (all sub-agent prompts live in their agent/ modules)
CHAT_SYSTEM_PROMPT = (
    "You are a conversational AI assistant, proficient in delivering high-quality "
    "responses and resolving tasks effectively. You are very attentive and respond "
    "in markdown format."
)
DOC_CHAT_SYSTEM_PROMPT = (
    "You are a conversational assistant, expert at providing quality and accurate "
    "answers based on a document(s) and/or image(s) provided. You are very attentive "
    "and respond in markdown format. Take your time to read through the document(s) "
    "and/or image(s) carefully and pay attention to relevant areas pertaining to the "
    "question(s). Once done reading, provide an answer to the user question(s)."
)

if 'messages' not in st.session_state:
    st.session_state['messages'] = []
if 'input_token' not in st.session_state:
    st.session_state['input_token'] = 0
if 'output_token' not in st.session_state:
    st.session_state['output_token'] = 0
if 'user_sess' not in st.session_state:
    st.session_state['user_sess'] =str(time.time())
if 'chat_session_list' not in st.session_state:
    st.session_state['chat_session_list'] = []
if 'count' not in st.session_state:
    st.session_state['count'] = 0
if 'userid' not in st.session_state:
    st.session_state['userid']= config_file["UserId"]
if 'cost' not in st.session_state:
    st.session_state['cost'] = 0
if 'reasoning_mode' not in st.session_state:
    st.session_state['reasoning_mode'] = False  # Only activated when user selects anthropic 3.7 and toggles on thinking

def compute_turn_cost(usage, prices, engine="runtime"):
    """Dollar cost of one turn from strands' accumulated usage, cache-aware.

    Only caching-capable models emit cacheReadInputTokens/cacheWriteInputTokens in
    the strands usage object (the keys are absent otherwise), so cache pricing
    applies exactly when the model reported cache traffic. Semantics differ by
    engine and MUST NOT be mixed up:
      - Bedrock Converse ("runtime"): cache tokens are reported SEPARATELY from
        inputTokens -> all four terms are additive.
      - OpenAI-style (Mantle engines): cached tokens are a SUBSET of inputTokens ->
        carve them out of the input term before pricing, else they bill twice.
    Models whose pricing entry lacks cache rates fall back to the input rate, so
    reported cache traffic is never silently free.
    """
    input_tokens = usage.get("inputTokens", 0)
    output_tokens = usage.get("outputTokens", 0)
    cache_read = usage.get("cacheReadInputTokens", 0)
    cache_write = usage.get("cacheWriteInputTokens", 0)

    if engine != "runtime" and cache_read:
        input_tokens = max(0, input_tokens - cache_read)

    cost = input_tokens * prices["input"] + output_tokens * prices["output"]
    if cache_read:
        cost += cache_read * prices.get("cache_read", prices["input"])
    if cache_write:
        cost += cache_write * prices.get("cache_write", prices["input"])
    return cost


def presign_s3_uri(s3_uri, expires_in=3600):
    """Presigned GET url for an s3:// uri. Minted at render time on every Streamlit
    rerun, so displayed links are always freshly signed (no refresh button needed).
    Note the effective lifetime is capped by the signing credentials' expiry."""
    bucket_name, key = s3_uri.replace('s3://', '').split('/', 1)
    return boto3.client('s3', region_name=REGION, config=config).generate_presigned_url(
        'get_object', Params={'Bucket': bucket_name, 'Key': key}, ExpiresIn=expires_in)


def render_artifacts_expander(doc_uris):
    """The 'artifacts' expander under an assistant message: one presigned download
    link per generated document."""
    if not doc_uris:
        return
    with st.expander(label="**artifacts**"):
        for uri in doc_uris:
            try:
                st.markdown(f"📄 [{os.path.basename(uri)}]({presign_s3_uri(uri)})")
            except Exception:
                st.markdown(f"📄 {os.path.basename(uri)} (link unavailable)")


def get_object_with_retry(bucket, key):
    max_retries = 5
    retries = 0
    backoff_base = 2
    max_backoff = 3  # Maximum backoff time in seconds
    s3 = boto3.client('s3')
    while retries < max_retries:
        try:
            response = s3.get_object(Bucket=bucket, Key=key)
            return response
        except ClientError as e:
            error_code = e.response['Error']['Code']
            if error_code == 'DecryptionFailureException':
                sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                print(f"Decryption failed, retrying in {sleep_time} seconds...")
                time.sleep(sleep_time)
                retries += 1
            elif e.response['Error']['Code'] == 'ModelStreamErrorException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
            else:
                raise e

    # If we reach this point, it means the maximum number of retries has been exceeded
    raise Exception(f"Failed to get object {key} from bucket {bucket} after {max_retries} retries.")
def process_files(files, cutoff=None):
    """process uploaded files in parallel.

    cutoff: max rows of tabular files (csv/xlsx) to inject into the conversation;
        None injects everything. Used when the Advanced Data Analytics tool is
        active - the sub-agent reads the FULL file from S3, so the conversation
        only needs a schema preview."""
    result_string=""
    errors = []
    future_proxy_mapping = {} 
    futures = []

    with concurrent.futures.ProcessPoolExecutor() as executor:
        # Partial function to pass the handle_doc_upload_or_s3 function
        func = partial(handle_doc_upload_or_s3, cutoff=cutoff)   
        for file in files:
            future = executor.submit(func, file)
            future_proxy_mapping[future] = file
            futures.append(future)

        # Collect the results and handle exceptions
        for future in concurrent.futures.as_completed(futures):        
            file_url= future_proxy_mapping[future]
            try:
                result = future.result()               
                doc_name=os.path.basename(file_url)                
                result_string+=f"<{doc_name}>\n{result}\n</{doc_name}>\n" # tag documnets with names to enhance prompts
            except Exception as e:
                # Get the original function arguments from the Future object
                error = {'file': file_url, 'error': str(e)}
                errors.append(error)

    return errors, result_string

def handle_doc_upload_or_s3(file, cutoff=None):
    """Handle various document format"""
    dir_name, ext = os.path.splitext(file)
    if  ext.lower() in [".pdf", ".png", ".jpg",".tif",".jpeg"]:   
        content=exract_pdf_text_aws(file)
    elif ".csv"  == ext.lower():
        content=parse_csv_from_s3(file,cutoff)
    elif ext.lower() in [".xlsx", ".xls"]:
        content=table_parser_utills(file,cutoff)   
    elif  ".json"==ext.lower():      
        obj=get_s3_obj_from_bucket_(file)
        content = json.loads(obj['Body'].read())  
    elif  ext.lower() in [".txt", ".py", ".md"]:       
        obj=get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
    elif ".docx" == ext.lower():       
        obj=get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        docx_buffer = io.BytesIO(content)
        content = extract_text_and_tables(docx_buffer)
    elif ".pptx" == ext.lower():       
        obj=get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        docx_buffer = io.BytesIO(content)        
        content = extract_text_from_pptx_s3(docx_buffer)
    # Implement any other file extension logic 
    return content

class InvalidContentError(Exception):
    pass

def detect_encoding(s3_uri):
    """detect csv encoding"""
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", s3_uri)
    if match:
        bucket_name = match.group(1)
        key = match.group(2) 
    response = s3.get_object(Bucket=bucket_name, Key=key)
    content = response['Body'].read()
    result = chardet.detect(content)
    return result['encoding']

def parse_csv_from_s3(s3_uri, cutoff):
    """read csv files"""
    try:
        # Detect the file encoding using chardet
        encoding = detect_encoding(s3_uri)        
        # Sniff the delimiter and read the CSV file
        df = pd.read_csv(s3_uri, delimiter=None, engine='python', encoding=encoding)
        if cutoff and len(df) > cutoff:
            total = len(df)
            return (df.iloc[:cutoff].to_csv(index=False, sep=CSV_SEPERATOR)
                    + f"\n[PREVIEW: first {cutoff} of {total} rows. The data-analysis "
                    "tool has access to the complete file.]\n")
        return df.to_csv(index=False, sep=CSV_SEPERATOR)
    except Exception as e:
        raise InvalidContentError(f"Error: {e}")
    
def iter_block_items(parent):
    if isinstance(parent, Document):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError("something's not right")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield DocxTable(child, parent)

def extract_text_and_tables(docx_path):
    """ Extract text from docx files"""
    document = DocxDocument(docx_path)
    content = ""
    current_section = ""
    section_type = None
    for block in iter_block_items(document):
        if isinstance(block, Paragraph):
            if block.text:
                if block.style.name == 'Heading 1':
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                        section_type = None  
                    section_type ="h1"
                    content += f"<{section_type}>{block.text}</{section_type}>\n"
                elif block.style.name== 'Heading 3':
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                    section_type = "h3"  
                    content += f"<{section_type}>{block.text}</{section_type}>\n"                
                elif block.style.name == 'List Paragraph':
                    # Add to the current list section
                    if section_type != "list":
                        # Close the current section if it exists
                        if current_section:
                            content += f"{current_section}</{section_type}>\n"
                        section_type = "list"
                        current_section = "<list>"
                    current_section += f"{block.text}\n"
                elif block.style.name.startswith('toc'):
                    # Add to the current toc section
                    if section_type != "toc":
                        # Close the current section if it exists
                        if current_section:
                            content += f"{current_section}</{section_type}>\n"
                        section_type = "toc"
                        current_section = "<toc>"
                    current_section += f"{block.text}\n"
                else:
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                        section_type = None

                    # Append the passage text without tagging
                    content += f"{block.text}\n"
        
        elif isinstance(block, DocxTable):
            # Add the current section before the table
            if current_section:
                content += f"{current_section}</{section_type}>\n"
                current_section = ""
                section_type = None

            content += "<table>\n"
            for row in block.rows:
                row_content = []
                for cell in row.cells:
                    cell_content = []
                    for nested_block in iter_block_items(cell):
                        if isinstance(nested_block, Paragraph):
                            cell_content.append(nested_block.text)
                        elif isinstance(nested_block, DocxTable):
                            nested_table_content = parse_nested_table(nested_block)
                            cell_content.append(nested_table_content)
                    row_content.append(CSV_SEPERATOR.join(cell_content))
                content += CSV_SEPERATOR.join(row_content) + "\n"
            content += "</table>\n"

    # Add the final section
    if current_section:
        content += f"{current_section}</{section_type}>\n"
    return content

def parse_nested_table(table):
    nested_table_content = "<table>\n"
    for row in table.rows:
        row_content = []
        for cell in row.cells:
            cell_content = []
            for nested_block in iter_block_items(cell):
                if isinstance(nested_block, Paragraph):
                    cell_content.append(nested_block.text)
                elif isinstance(nested_block, DocxTable):
                    nested_table_content += parse_nested_table(nested_block)
            row_content.append(CSV_SEPERATOR.join(cell_content))
        nested_table_content += CSV_SEPERATOR.join(row_content) + "\n"
    nested_table_content += "</table>"
    return nested_table_content



def extract_text_from_pptx_s3(pptx_buffer):
    """ Extract Text from pptx files"""
    presentation = Presentation(pptx_buffer)    
    text_content = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text)
        text_content.append('\n'.join(slide_text))    
    return '\n\n'.join(text_content)
    
def exract_pdf_text_aws(file):
    """extract text from PDFs using Amazon Textract or PyPDF2"""
    file_base_name = os.path.basename(file)
    dir_name, ext = os.path.splitext(file)
    # Checking if extracted doc content is in S3
    if USE_TEXTRACT:        
        if [x for x in get_s3_keys(f"{TEXTRACT_RESULT_CACHE_PATH}/") if file_base_name in x]:    
            response = get_object_with_retry(BUCKET, f"{TEXTRACT_RESULT_CACHE_PATH}/{file_base_name}.txt")
            text = response['Body'].read().decode()
            return text
        else:
            
            extractor = Textractor()
            # Asynchronous call, you will experience some wait time. Try caching results for better experience
            if "pdf" in ext:
                print("Asynchronous call, you may experience some wait time.")
                document = extractor.start_document_analysis(
                    file_source=file,
                    features=[TextractFeatures.LAYOUT, TextractFeatures.TABLES],
                    save_image=False,
                    s3_output_path=f"s3://{BUCKET}/textract_output/"
                )
            # Synchronous call
            else:
                document = extractor.analyze_document(
                    file_source=file,
                    features=[TextractFeatures.LAYOUT,TextractFeatures.TABLES],  
                    save_image=False,
                  )
            config = TextLinearizationConfig(
                hide_figure_layout=False,   
                hide_header_layout=False,    
                table_prefix="<table>",
                table_suffix="</table>",
            )
            # Upload extracted content to s3
            S3.put_object(Body=document.get_text(config=config), Bucket=BUCKET, Key=f"{TEXTRACT_RESULT_CACHE_PATH}/{file_base_name}.txt") 
            return document.get_text(config=config)
    else:
        s3 = boto3.resource("s3")
        match = re.match("s3://(.+?)/(.+)", file)
        if match:
            bucket_name = match.group(1)
            key = match.group(2)
        if "pdf" in ext:            
            pdf_bytes = io.BytesIO()            
            s3.Bucket(bucket_name).download_fileobj(key, pdf_bytes)
            # Read the PDF from the BytesIO object
            pdf_bytes.seek(0)                      
            # Create a PDF reader object
            pdf_reader = PyPDF2.PdfReader(pdf_bytes)
            # Get the number of pages in the PDF
            num_pages = len(pdf_reader.pages)
            # Extract text from each page
            text = ''
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text += page.extract_text()
        else:
            img_bytes = io.BytesIO()
            s3.Bucket(bucket_name).download_fileobj(key, img_bytes)
            img_bytes.seek(0)         
            image_stream = io.BytesIO(img_bytes)
            image = Image.open(image_stream)
            text = pytesseract.image_to_string(image)
        return text

def strip_newline(cell):
    return str(cell).strip()

def table_parser_openpyxl(file, cutoff):
    """convert xlsx files to python pandas handling merged cells"""
    # Read from S3
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
        obj = s3.get_object(Bucket=bucket_name, Key=key)    
        # Read Excel file from S3 into a buffer
        xlsx_buffer = io.BytesIO(obj['Body'].read())
        xlsx_buffer.seek(0)    
        # Load workbook
        wb = openpyxl.load_workbook(xlsx_buffer)    
        all_sheets_string = ""
        # Iterate over each sheet in the workbook
        for sheet_name in wb.sheetnames:
            # all_sheets_name.append(sheet_name)
            worksheet = wb[sheet_name]

            all_merged_cell_ranges: list[CellRange] = list(
                worksheet.merged_cells.ranges
            )
            for merged_cell_range in all_merged_cell_ranges:
                merged_cell: Cell = merged_cell_range.start_cell
                worksheet.unmerge_cells(range_string=merged_cell_range.coord)
                for row_index, col_index in merged_cell_range.cells:
                    cell: Cell = worksheet.cell(row=row_index, column=col_index)
                    cell.value = merged_cell.value        
            # Convert sheet data to a DataFrame
            df = pd.DataFrame(worksheet.values)
            df = df.map(strip_newline)
            preview_note = ""
            if cutoff and len(df) > cutoff:
                preview_note = (f"\n[PREVIEW: first {cutoff} of {len(df)} rows. The "
                                "data-analysis tool has access to the complete file.]\n")
                df = df.iloc[:cutoff]

            # Convert to string and tag by sheet name
            tabb=df.to_csv(sep=CSV_SEPERATOR, index=False, header=0)
            all_sheets_string+=f'<{sheet_name}>\n{tabb}{preview_note}\n</{sheet_name}>\n'
        return all_sheets_string
    else:
        raise Exception(f"{file} not formatted as an S3 path")

def calamaine_excel_engine(file,cutoff):
    # # Read from S3
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
        obj = s3.get_object(Bucket=bucket_name, Key=key)    
        # Read Excel file from S3 into a buffer
        xlsx_buffer = io.BytesIO(obj['Body'].read())
        xlsx_buffer.seek(0)    
        all_sheets_string = ""
        # Load the Excel file
        workbook = CalamineWorkbook.from_filelike(xlsx_buffer)
        # Iterate over each sheet in the workbook
        for sheet_name in workbook.sheet_names:
            # Get the sheet by name
            sheet = workbook.get_sheet_by_name(sheet_name)
            df = pd.DataFrame(sheet.to_python(skip_empty_area=False))
            df = df.map(strip_newline)
            preview_note = ""
            if cutoff and len(df) > cutoff:
                preview_note = (f"\n[PREVIEW: first {cutoff} of {len(df)} rows. The "
                                "data-analysis tool has access to the complete file.]\n")
                df = df.iloc[:cutoff]
            # print(df)
            tabb = df.to_csv(sep=CSV_SEPERATOR, index=False, header=0)
            all_sheets_string += f'<{sheet_name}>\n{tabb}{preview_note}\n</{sheet_name}>\n'
        return all_sheets_string
    else:
        raise Exception(f"{file} not formatted as an S3 path")

def table_parser_utills(file,cutoff):
    try:
        response = table_parser_openpyxl(file, cutoff)
        if response:
            return response
        else:
            return calamaine_excel_engine(file, cutoff)        
    except Exception as e:
        try:
            return calamaine_excel_engine(file, cutoff)
        except Exception as e:
            raise Exception(str(e))
def get_s3_keys(prefix):
    """list all keys in an s3 path"""
    s3 = boto3.client('s3')
    keys = []
    next_token = None
    while True:
        if next_token:
            response = s3.list_objects_v2(Bucket=BUCKET, Prefix=prefix, ContinuationToken=next_token)
        else:
            response = s3.list_objects_v2(Bucket=BUCKET, Prefix=prefix)
        if "Contents" in response:
            for obj in response['Contents']:
                key = obj['Key']
                name = key[len(prefix):]
                keys.append(name)
        if "NextContinuationToken" in response:
            next_token = response["NextContinuationToken"]
        else:
            break
    return keys
    
def parse_s3_uri(uri):
    """
    Parse an S3 URI and extract the bucket name and key.

    :param uri: S3 URI (e.g., 's3://bucket-name/path/to/file.txt')
    :return: Tuple of (bucket_name, key) if valid, (None, None) if invalid
    """
    pattern = r'^s3://([^/]+)/(.*)$'
    match = re.match(pattern, uri)
    if match:
        return match.groups()
    return (None, None)
    
def copy_s3_object(source_uri, dest_bucket, dest_key):
    """
    Copy an object from one S3 location to another.

    :param source_uri: S3 URI of the source object
    :param dest_bucket: Name of the destination bucket
    :param dest_key: Key to be used for the destination object
    :return: True if successful, False otherwise
    """
    s3 = boto3.client('s3')

    # Parse the source URI
    source_bucket, source_key = parse_s3_uri(source_uri)
    if not source_bucket or not source_key:
        print(f"Invalid source URI: {source_uri}")
        return False

    try:
        # Create a copy source dictionary
        copy_source = {
            'Bucket': source_bucket,
            'Key': source_key
        }
        destination_key = f"{dest_key}/{os.path.basename(source_key)}"
        # Copy the object
        s3.copy_object(CopySource=copy_source, Bucket=dest_bucket, Key=destination_key)

        print(f"File copied from {source_uri} to s3://{dest_bucket}/{destination_key}")
        return f"s3://{dest_bucket}/{destination_key}"

    except ClientError as e:
        print(f"An error occurred: {e}")
        raise(e)
        # return False

def plotly_to_png_bytes(s3_uri):
    """
    Read a .plotly file from S3 given an S3 URI, convert it to a PNG image, and return the image as bytes.

    :param s3_uri: S3 URI of the .plotly file (e.g., 's3://bucket-name/path/to/file.plotly')
    :return: PNG image as bytes
    """
    # Parse S3 URI
    parsed_uri = urlparse(s3_uri)
    bucket_name = parsed_uri.netloc
    file_key = parsed_uri.path.lstrip('/')

    # Initialize S3 client
    s3_client = boto3.client('s3')

    try:
        # Read the .plotly file from S3
        response = s3_client.get_object(Bucket=bucket_name, Key=file_key)
        plotly_data = json.loads(response['Body'].read().decode('utf-8'))

        # Create a Figure object from the plotly data
        fig = go.Figure(data=plotly_data['data'], layout=plotly_data.get('layout', {}))

        # Convert the figure to PNG bytes
        img_bytes = fig.to_image(format="png")

        return img_bytes

    except Exception as e:
        print(f"An error occurred: {str(e)}")
        return None
        

def get_s3_obj_from_bucket_(file):
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)    
        obj = s3.get_object(Bucket=bucket_name, Key=key)  
    return obj

def put_obj_in_s3_bucket_(docs, key_prefix=S3_DOC_CACHE_PATH):
    """Cache an uploaded file (or copy an existing s3 object) under key_prefix.

    key_prefix is session-scoped by the caller (uploads/<session_id>) so files with
    the same name from different sessions don't overwrite each other in the cache.
    """
    if isinstance(docs,str):
        s3_uri_pattern = r'^s3://([^/]+)/(.*?([^/]+)/?)$'
        if bool(re.match(s3_uri_pattern,  docs)):
            file_uri=copy_s3_object(docs, BUCKET, key_prefix)
            return file_uri
    else:
        file_name = os.path.basename(docs.name)
        file_path = f"{key_prefix}/{file_name}"
        S3.put_object(Body=docs.read(), Bucket= BUCKET, Key=file_path)
        return f"s3://{BUCKET}/{file_path}"


def image_blocks_from_s3(image_paths):
    """Load s3 images (or plotly json rendered to png) as Converse image content blocks"""
    content = []
    s3 = boto3.client('s3', region_name=REGION)
    for img in image_paths:
        match = re.match("s3://(.+?)/(.+)", img)
        image_name = os.path.basename(img)
        _, ext = os.path.splitext(image_name)
        if "jpg" in ext:
            ext = ".jpeg"
        bucket_name = match.group(1)
        key = match.group(2)
        if ".plotly" in key:
            bytes_image = plotly_to_png_bytes(img)
            ext = ".png"
        else:
            obj = s3.get_object(Bucket=bucket_name, Key=key)
            bytes_image = obj['Body'].read()
        content.extend([{"text": image_name}, {
            "image": {
                "format": f"{ext.lower().replace('.', '')}",
                "source": {"bytes": bytes_image}
            }
        }])
    return content
def get_session_ids_by_user(table_name, user_id):
    """
    Get Session Ids and corresponding top message for a user to populate the chat
    history drop down on the front end. Reconstructed from the configured Strands
    session backend (local | s3 | dynamodb).
    """
    return list_user_sessions(
        SESSION_STORAGE, user_id=user_id, bucket=BUCKET, region=REGION,
        dynamodb_table=table_name, agentcore_memory_id=AGENTCORE_MEMORY_ID,
    )

def list_csv_xlsx_in_s3_folder(bucket_name, folder_path):
    """
    List all CSV and XLSX files in a specified S3 folder.

    :param bucket_name: Name of the S3 bucket
    :param folder_path: Path to the folder in the S3 bucket
    :return: List of CSV and XLSX file names in the folder
    """
    s3 = boto3.client('s3')
    csv_xlsx_files = []

    try:
        # Ensure the folder path ends with a '/'
        if not folder_path.endswith('/'):
            folder_path += '/'

        # List objects in the specified folder
        paginator = s3.get_paginator('list_objects_v2')
        page_iterator = paginator.paginate(Bucket=bucket_name, Prefix=folder_path)

        for page in page_iterator:
            if 'Contents' in page:
                for obj in page['Contents']:
                    # Get the file name
                    file_name = obj['Key']

                    # Check if the file is a CSV or XLSX
                    if file_name.lower().endswith(INPUT_EXT):
                        csv_xlsx_files.append(os.path.basename(file_name))
                        # csv_xlsx_files.append(file_name)

        return csv_xlsx_files

    except ClientError as e:
        print(f"An error occurred: {e}")
        return []

def query_llm(params, handler, workings_slot=None):
    """
    Handles users requests and routes to a native call or tool use, then stores sonversation to local or DynamoDB
    """  

    if not isinstance(params['upload_doc'], list):
        raise TypeError("documents must be in a list format")

    vision_model = True
    # display name goes straight to build_chat_agent; resolve_model routes it to the
    # right provider (runtime Converse vs Mantle OpenAI) using the registry spec
    model = params['model']
    if any(keyword in [params['model']] for keyword in NON_VISION_MODELS):
        vision_model = False

    # Streaming/activity renderer for the whole turn. Created BEFORE attachment
    # ingestion so uploads and text extraction (Textract can take a while) show a
    # live activity line instead of a blank bubble; the line is replaced as soon as
    # model output starts streaming, same as the tool spinners.
    stream_handler = StreamlitCallbackHandler(handler)

    # prompt template for when a user uploads a doc
    doc_path = []
    image_path = []
    full_doc_path = []
    doc = ""
    if params['upload_doc'] or params['s3_objects']:
        # session-scope the upload cache so same-named files from different
        # sessions don't collide in s3
        upload_prefix = f"{S3_DOC_CACHE_PATH}/{params['session_id']}"
        n_files = len(params['upload_doc'] or []) + len(params['s3_objects'] or [])
        uploaded = 0
        if params['upload_doc']:
            for ids, docs in enumerate(params['upload_doc']):
                file_name = docs.name
                uploaded += 1
                stream_handler.status(f"📤 Uploading {file_name} ({uploaded}/{n_files})…")
                _, extensions = os.path.splitext(file_name)
                docs = put_obj_in_s3_bucket_(docs, upload_prefix)
                full_doc_path.append(docs)
                if extensions.lower() in [".jpg", ".jpeg", ".png", ".gif", ".webp"] and vision_model:
                    image_path.append(docs)
                    continue

        if params['s3_objects']:
            for ids, docs in enumerate(params['s3_objects']):
                file_name = docs
                uploaded += 1
                stream_handler.status(f"📤 Staging {file_name} ({uploaded}/{n_files})…")
                _, extensions = os.path.splitext(file_name)
                docs = put_obj_in_s3_bucket_(f"s3://{INPUT_BUCKET}/{INPUT_S3_PATH}/{docs}", upload_prefix)
                full_doc_path.append(docs)
                if extensions.lower() in [".jpg", ".jpeg", ".png", ".gif", ".webp"] and vision_model:
                    image_path.append(docs)
                    continue

        doc_path = [item for item in full_doc_path if item not in image_path]
        if doc_path:
            names = ", ".join(os.path.basename(p) for p in doc_path)
            stream_handler.status(f"📄 Processing attachment(s): {names[:120]}…")
        # With the data-analysis tool active, inject only a 10-row preview of tabular
        # files (csv/xlsx) - enough schema/flavor for the orchestrator to converse and
        # write good analysis briefs, while the sub-agent reads the FULL file from S3
        # in its sandbox. Non-tabular formats are always injected in full; without the
        # tool, full injection is the only way the model sees the data at all.
        tabular_cutoff = 10 if "Advanced Data Analytics" in params['tools'] else None
        errors, result_string = process_files(doc_path, cutoff=tabular_cutoff)
        if errors:
            st.error(errors)
        # Wrap injected docs in a sentinel so the clean question can be recovered
        # from the persisted message for display (see agent/sessions.py).
        doc = wrap_attached_documents(result_string)
        chat_template = DOC_CHAT_SYSTEM_PROMPT
        # When tabular attachments were truncated to previews, tell the orchestrator
        # up front in the system prompt: previews are for schema/context only - any
        # computation over the data must go through the data_analysis tool.
        if tabular_cutoff and "[PREVIEW: first" in result_string:
            chat_template += (
                "\n\nNote: attached tabular files (CSV/Excel) are shown as PREVIEWS "
                f"(first {tabular_cutoff} rows) marked with a [PREVIEW: ...] line - "
                "the full datasets are larger. Use the previews only to understand "
                "schema and content. For ANY computation, aggregation, filtering or "
                "chart over the data, use the data_analysis tool, which reads the "
                "complete file; never compute answers from the preview rows alone."
            )
    else:
        chat_template = CHAT_SYSTEM_PROMPT

    # Preprocessed docs are injected into the user turn; images ride as content blocks
    content = image_blocks_from_s3(image_path) if image_path else []
    content.append({"text": doc + params['question']})

    # Enable web tools for this turn if the user selected "Web Search" in the sidebar
    agent_tools = list(WEB_TOOLS) if "Web Search" in params['tools'] else []

    # Worker-agent swarm: parallel ephemeral web-research agents the orchestrator
    # can delegate self-contained subtasks to. Stateless; progress is line-per-worker
    # in the activity spinner (no per-worker output streaming - it would interleave).
    if "Research Workers" in params['tools']:
        agent_tools.append(make_worker_agents_tool(
            model_id=WORKER_AGENT_MODEL,
            region=REGION,
            max_workers=WORKER_AGENT_MAX_WORKERS,
            status_fn=stream_handler.status,
        ))

    # Data-analysis agent-as-tool: an ephemeral sub-agent that writes & runs code in a
    # sandbox picked by the sidebar Runtime slider — "python" = AgentCore Code
    # Interpreter, "pyspark" = Athena for Apache Spark (PySpark engine v3). The
    # execution session and the sub-agent message history are cached per (chat
    # session, engine) so dataframes/context survive turns but never leak across
    # engines (the two sandboxes have different environments and prompts).
    da_sink = {"image_output": [], "plotly": [], "doc_output": []}
    # Datasets offered to the DA/docgen sub-agents: this turn's attachments now,
    # extended IN PLACE after the agent is built with previously attached docs
    # still visible in the conversation window (see dataset_uris_in_window). The
    # tool closures hold this list by reference and read it at invocation time.
    da_datasets = list(doc_path)
    if "Advanced Data Analytics" in params['tools']:
        da_engine = params.get("engine") or "python"
        da_prefix = f"{S3_DOC_CACHE_PATH}/{params['session_id']}"
        if da_engine == "pyspark" and not ATHENA_WORKGROUP:
            st.error("The pyspark runtime requires 'athena-work-group-name' in config.json")
        elif da_engine != "pyspark" and not CODE_INTERPRETER_ID:
            st.error("Advanced Data Analytics requires 'code-interpreter-id' in config.json")
        else:
            da_key = f"{params['session_id']}:{da_engine}"
            da_state = st.session_state.setdefault('data_analysis', {})
            entry = da_state.get(da_key)
            if entry is None:
                if da_engine == "pyspark":
                    da_session = AthenaSparkSessionManager(
                        REGION, ATHENA_WORKGROUP, bucket=BUCKET, prefix=da_prefix)
                else:
                    da_session = CodeSessionManager(
                        REGION, CODE_INTERPRETER_ID,
                        timeout_seconds=CODE_INTERPRETER_TIMEOUT)
                entry = {"session": da_session, "store": {"messages": []}}
                da_state[da_key] = entry
            da_model_id = DATA_ANALYSIS_MODEL
            da_vision = not any(k in [DATA_ANALYSIS_MODEL] for k in NON_VISION_MODELS)
            da_panel = (SubAgentWorkingsPanel(workings_slot, "🔬 Agent workings")
                        if workings_slot is not None else None)
            agent_tools.append(make_data_analysis_tool(
                dataset_uris=da_datasets,
                model_id=da_model_id,
                region=REGION,
                bucket=BUCKET,
                upload_prefix=da_prefix,
                session=entry["session"],
                artifact_sink=da_sink,
                message_store=entry["store"],
                vision=da_vision,
                status_fn=stream_handler.status,
                workings_panel=da_panel,
            ))

    # Document-generator agent-as-tool: its own interpreter session + message store
    # (cached under the "docgen" key), sharing the DA model and the per-turn artifact
    # sink - generated files land in doc_output and render in the artifacts expander.
    dg_generated = []  # generated-artifact registry; filled after the agent is built
    if "Document Generator" in params['tools']:
        if not CODE_INTERPRETER_ID:
            st.error("Document Generator requires 'code-interpreter-id' in config.json")
        else:
            # artifacts get their own prefix so an edit saved under an attached
            # file's basename can never overwrite the cached original attachment
            dg_prefix = f"{S3_DOC_CACHE_PATH}/{params['session_id']}/artifacts"
            dg_key = f"{params['session_id']}:docgen"
            da_state = st.session_state.setdefault('data_analysis', {})
            dg_entry = da_state.get(dg_key)
            if dg_entry is None:
                dg_entry = {
                    "session": CodeSessionManager(REGION, CODE_INTERPRETER_ID,
                                                  timeout_seconds=CODE_INTERPRETER_TIMEOUT),
                    "store": {"messages": []},
                }
                da_state[dg_key] = dg_entry
            dg_model_id = DATA_ANALYSIS_MODEL
            dg_panel = (SubAgentWorkingsPanel(workings_slot, "📄 Agent workings")
                        if workings_slot is not None else None)
            agent_tools.append(make_document_generator_tool(
                model_id=dg_model_id,
                region=REGION,
                bucket=BUCKET,
                upload_prefix=dg_prefix,
                available_uris=da_datasets,
                generated_uris=dg_generated,
                session=dg_entry["session"],
                artifact_sink=da_sink,
                message_store=dg_entry["store"],
                status_fn=stream_handler.status,
                workings_panel=dg_panel,
            ))
    chat_agent = build_chat_agent(
        model_id=model,
        region=REGION,
        session_id=params["session_id"],
        session_storage=SESSION_STORAGE,
        system_prompt=chat_template,
        history_window=CHAT_HISTORY_LENGTH * 2,  # config counts QA pairs; window counts messages
        reasoning=st.session_state['reasoning_mode'],
        max_tokens=OUTPUT_TOKEN,
        user_id=st.session_state['userid'],
        bucket=BUCKET,
        dynamodb_table=DYNAMODB_TABLE,
        agentcore_memory_id=AGENTCORE_MEMORY_ID,
        tools=agent_tools,
        callback_handler=stream_handler,
    )
    # Extend the sub-agents' dataset list with previously attached docs whose
    # content is still visible in the restored window (turn_meta holds their uris).
    # In-place so the already-built tool closures see the additions.
    prior_meta = chat_agent.state.get("turn_meta") or {}
    for uri in dataset_uris_in_window(chat_agent.messages, prior_meta):
        if uri not in da_datasets and uri.lower().endswith(INPUT_EXT):
            da_datasets.append(uri)
    # ...and the docgen registry with generated artifacts still referenced in the
    # window (their [generated artifacts: ...] markers), so edit requests survive
    # restarts. In-place: the docgen closure holds dg_generated by reference.
    dg_generated.extend(generated_uris_in_window(chat_agent.messages, prior_meta))

    try:
        result = chat_agent(content)
    except Exception as e:
        # A failed/stalled model stream must not leave the UI stuck: stop late
        # worker-thread writes to the placeholder, surface the error, end the turn
        # cleanly. The turn's user message may already be persisted - harmless, the
        # next turn continues the conversation (verified).
        stream_handler.close()
        st.error(f"Model call failed: {type(e).__name__}: {e}")
        return ("Sorry - the model call failed before completing "
                f"({type(e).__name__}). Please try again.")
    finally:
        stream_handler.close()
    response = str(result)

    usage = result.metrics.accumulated_usage
    st.session_state['input_token'] = usage.get("inputTokens", 0)
    st.session_state['output_token'] = usage.get("outputTokens", 0)
    turn_cost = compute_turn_cost(
        usage, pricing_file[params['model']],
        engine=model_info[params['model']].get("engine", "runtime"))
    st.session_state['cost'] += turn_cost

    # Persist the un-reconstructable per-turn metadata (model id, cost, and the
    # s3 uris of attachments for the citations dropdown) in agent.state, keyed by
    # the assistant message id. Everything else (text, thinking, tokens, timestamps)
    # is read back from the persisted messages themselves.
    assistant_message_id = chat_agent._session_manager._latest_agent_message[chat_agent.agent_id].message_id
    turn_meta = chat_agent.state.get("turn_meta") or {}
    turn_meta[str(assistant_message_id)] = {
        "model": params['model'],
        "modelID": model,
        "cost": turn_cost,
        "documents": doc_path,
        "images": image_path,
        # data-analysis artifacts (chart/plotly s3 uris) produced this turn, if any
        "image_output": da_sink["image_output"],
        "plotly": da_sink["plotly"],
        "doc_output": da_sink["doc_output"],
    }
    chat_agent.state.set("turn_meta", turn_meta)
    chat_agent._session_manager.sync_agent(chat_agent)
    return response


def get_chat_historie_for_streamlit(params):
    """Reconstruct the display log for a session directly from its Strands session
    (persisted messages + per-turn metadata in agent.state), for the configured
    backend (local | s3 | dynamodb | agentcore). No separate display log is written."""
    return read_display_history(
        SESSION_STORAGE, params["session_id"],
        user_id=st.session_state['userid'], bucket=BUCKET, region=REGION,
        dynamodb_table=DYNAMODB_TABLE, agentcore_memory_id=AGENTCORE_MEMORY_ID,
    )


def get_key_from_value(dictionary, value):
    return next((key for key, val in dictionary.items() if val == value), None)


def escape_dollars_outside_code(text: str, replacement: str = "\\$") -> str:
    """Escape $ signs in markdown, but leave fenced code blocks (``` ... ```) untouched.
    Issue is that Streamlit struggles to render $ without proper escaping or in (```)
    encapsulation. The `chat_debrock_` function checks for (```) and renders the content
    as is, however, there may be $ characters outside the (```) in the message content
    that needs to be handled seperately"""
    # Pattern captures fenced code blocks (with optional language hint),
    # non-greedy so multiple blocks are handled separately.
    fence_pattern = re.compile(r"(```.*?```)", re.DOTALL)

    parts = fence_pattern.split(text)

    # split() with a capturing group returns alternating segments:
    # [outside, fenced, outside, fenced, outside, ...]
    # Even indexes = outside code blocks -> escape them
    # Odd indexes  = inside code blocks  -> leave as-is
    for i in range(0, len(parts), 2):
        if "\\$" not in parts[i]: # Skip if LLM response already escaped "$"
            parts[i] = parts[i].replace("$", replacement)

    return "".join(parts)

def chat_bedrock_(params):
    st.title('Chatty AI Assitant 🙂')
    if params["session_id"].strip():
        st.session_state.messages = get_chat_historie_for_streamlit(params)
    for x,message in enumerate(st.session_state.messages):
        with st.chat_message(message["role"]):
            if "```" in message["content"]:
                message['content'] = escape_dollars_outside_code(message['content'])
                st.markdown(message["content"], unsafe_allow_html=True)
            else:
                st.markdown(message["content"].replace("$", "\\$"), unsafe_allow_html=True)
            if message["role"] == "assistant":
                for item in message["plotly"]:
                    bucket_name, key = item.replace('s3://', '').split('/', 1)
                    image_bytes = get_object_with_retry(bucket_name, key)
                    fig = pio.from_json(image_bytes['Body'].read().decode('utf-8'))
                    st.plotly_chart(fig)
                # PNGs that are just static twins of a plotly figure above are skipped
                # (same basename stem); standalone PNGs (e.g. matplotlib) still render.
                plotly_stems = {os.path.splitext(os.path.basename(p))[0] for p in message["plotly"]}
                for item in message["image_output"]:
                    if os.path.splitext(os.path.basename(item))[0] in plotly_stems:
                        continue
                    bucket_name, key = item.replace('s3://', '').split('/', 1)
                    image_bytes = get_object_with_retry(bucket_name, key)
                    image = Image.open(io.BytesIO(image_bytes['Body'].read()))
                    st.image(image)
                render_artifacts_expander(message.get("doc_output", []))
                if message["attachment"]:
                    with st.expander(label="**attachments**"):
                        st.markdown(message["attachment"])
                        # st.markdown(message["image_output"])
                if message['code']:
                    with st.expander(label="**code snippet**"):
                        st.markdown(f'```python\n{message["code"]}', unsafe_allow_html=True)
                    with st.expander(label="**code result**"):
                        st.markdown(f'```python\n{message["code-result"]}', unsafe_allow_html=True)
                if message['thinking']:
                    with st.expander(label="**MODEL REASONING**"):
                        st.markdown(message["thinking"].replace("$", "\\$"), unsafe_allow_html=True)

    if prompt := st.chat_input("Whats up?"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt.replace("$", "\\$"), unsafe_allow_html=True )
        with st.chat_message("assistant"):
            message_placeholder = st.empty()
            # slot for the transient "Agent workings" panel (fills while a sub-agent
            # runs, under the answer spinner); must be created on the main thread.
            workings_slot = st.empty()
            params["question"] = prompt
            answer = query_llm(params, message_placeholder, workings_slot)
            message_placeholder.markdown(answer.replace("$", "\\$"), unsafe_allow_html=True )
            st.session_state.messages.append({"role": "assistant", "content": answer})
            # This turn consumed the staged attachments (their content is now in the
            # conversation history) - rotate the widget nonce so the rerun remounts
            # the upload box / Files picker empty. Only reached on success: a failed
            # turn leaves the files staged for retry.
            if params['upload_doc'] or params['s3_objects']:
                st.session_state['attach_nonce'] += 1
        st.rerun()

def format_chat_session(session_string):
    # Truncate to first 10 words or 50 chars, whichever comes first
    words = session_string.split()[:10]
    truncated = ' '.join(words)

    if len(truncated) > 150:
        return truncated[:150] + "..."
    elif len(session_string.split()) > 10:
        return truncated + "..."
    else:
        return truncated

def app_sidebar():
    with st.sidebar:
        st.metric(label="Bedrock Session Cost", value=f"${round(st.session_state['cost'], 2)}") 
        st.write("-----")
        button = st.button("New Chat", type="primary")
        models = MODEL_DISPLAY_NAME
        model = st.selectbox('**Model**', models)
        if any(keyword in [model] for keyword in HYBRID_MODELS):
            st.session_state['reasoning_mode'] = st.toggle("Reasoning Mode", value=False, key="thinking")
        else:
            st.session_state['reasoning_mode'] = False
        runtime = ""
        tools = ""
        user_sess_id = get_session_ids_by_user(DYNAMODB_TABLE, st.session_state['userid'])
        float_keys = {float(key): value for key, value in user_sess_id.items()}
        sorted_messages = sorted(float_keys.items(), reverse=True)
        if button:
            # Fresh session id, clean pane, and REMOUNT the session selectbox under a
            # new key so it defaults to index 0 ("New Chat"). Merely clearing its
            # session_state entry is not enough: the browser-side component keeps the
            # old selection and re-submits it on the next interaction, silently
            # mapping the new conversation back onto the old session.
            st.session_state['user_sess'] = str(time.time())
            st.session_state['chat_list_nonce'] = st.session_state.get('chat_list_nonce', 0) + 1
            st.session_state['messages'] = []
        # single "New Chat" entry for the pending session id; once its first message
        # is sent it appears in storage with a real title, which wins the dict merge
        sorted_messages.insert(0, (float(st.session_state['user_sess']), "New Chat"))
        st.session_state['chat_session_list'] = dict(sorted_messages)
        chat_items = st.selectbox(
                            "**Chat Sessions**",
                            st.session_state['chat_session_list'].values(),
                            format_func=format_chat_session,
                            key=f"chat_sessions_{st.session_state.setdefault('chat_list_nonce', 0)}"
                        )
        #st.selectbox("**Chat Sessions**", st.session_state['chat_session_list'].values(), key="chat_sessions")
        session_id = get_key_from_value(st.session_state['chat_session_list'], chat_items)
        if model not in NON_TOOL_SUPPORTING_MODELS:
            tools = st.multiselect("**Tools**",
                                   ["Advanced Data Analytics", "Web Search",
                                    "Research Workers", "Document Generator"],
                                   key="function_collen", default=None)
            if "Advanced Data Analytics" in tools:
                engines = ["pyspark", "python"]
                runtime = st.select_slider(
                                "Runtime", engines, value="python", key="enginees"
                            )
        # Attachment widgets are nonce-keyed: after a query consumes attachments, the
        # nonce increments and the rerun remounts them EMPTY (attachments are per-turn;
        # the conversation history carries the content for follow-ups, and sub-agents
        # get still-visible dataset uris via dataset_uris_in_window). Reruns from other
        # widget changes leave the nonce - and therefore staged files - untouched.
        nonce = st.session_state.setdefault('attach_nonce', 0)
        bucket_items = list_csv_xlsx_in_s3_folder(INPUT_BUCKET, INPUT_S3_PATH)
        bucket_objects = st.multiselect("**Files**", bucket_items,
                                        key=f"objector_{nonce}", default=None)
        file = st.file_uploader('Upload a document', accept_multiple_files=True,
                                key=f"upload_doc_{nonce}",
                                help="pdf,csv,txt,png,jpg,xlsx,json,py doc format supported")
        if file or bucket_objects:
            st.caption("📎 Attachments are sent with your next question only, then cleared "
                       "— follow-ups answer from the conversation.")
        params = {"model": model, "session_id": str(session_id),
                  "chat_item": chat_items,
                  "upload_doc": file,
                  "tools": tools,
                  's3_objects': bucket_objects,
                  "engine": runtime
                 }
        st.session_state['count'] = 1
        return params


def main():
    params = app_sidebar()
    chat_bedrock_(params)
if __name__ == '__main__':
    main()
