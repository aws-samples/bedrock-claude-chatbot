import streamlit as st
import boto3
from botocore.config import Config
import os
import pandas as pd
import time
import json
from botocore.exceptions import ClientError
import io
import re
from pptx import Presentation
import random
from python_calamine import CalamineWorkbook
import chardet
from docx.table import _Cell
import concurrent.futures
from functools import partial
from textractor import Textractor
from textractor.data.constants import TextractFeatures
from textractor.data.text_linearization_config import TextLinearizationConfig
import pytesseract
from PIL import Image
import PyPDF2
from docx import Document as DocxDocument
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.document import Document
from docx.text.paragraph import Paragraph
from docx.table import Table as DocxTable
from utils.athena_handler_ import send_athena_job
import ast
from urllib.parse import urlparse
import plotly.graph_objects as go

config = Config(
    read_timeout=600,  # Read timeout parameter
    connect_timeout=600,  # Connection timeout parameter in seconds
    retries=dict(
        max_attempts=10  # Handle retries
    )
)

with open('config.json', 'r', encoding='utf-8') as f:
    config_file = json.load(f)

# Bedrock Model info
with open('model_id.json', 'r', encoding='utf-8') as f:
    model_info = json.load(f)

S3 = boto3.client('s3')
DYNAMODB = boto3.resource('dynamodb')
LOCAL_CHAT_FILE_NAME = "chat-history.json"
DYNAMODB_TABLE = config_file["DynamodbTable"]
BUCKET = config_file["Bucket_Name"]
OUTPUT_TOKEN = config_file["max-output-token"]
S3_DOC_CACHE_PATH = config_file["document-upload-cache-s3-path"]
TEXTRACT_RESULT_CACHE_PATH = config_file["AmazonTextract-result-cache"]
LOAD_DOC_IN_ALL_CHAT_CONVO = config_file["load-doc-in-chat-history"]
CHAT_HISTORY_LENGTH = config_file["chat-history-loaded-length"]
DYNAMODB_USER = config_file["UserId"]
REGION = config_file["region"]
USE_TEXTRACT = config_file["AmazonTextract"]
CSV_SEPERATOR = config_file["csv-delimiter"]
LAMBDA_FUNC = config_file["lambda-function"]
INPUT_S3_PATH = config_file["input_s3_path"]
INPUT_BUCKET = config_file["input_bucket"]
ATHENA_WORKGROUP_NAME = config_file["athena-work-group-name"]
TEXT_ONLY_MODELS = ["deepseek", "haiku-3.5", "micro"]

with open('pricing.json', 'r', encoding='utf-8') as f:
    pricing_file = json.load(f)



def put_db(params, messages):
    """Store long term chat history in DynamoDB"""    
    chat_item = {
        "UserId": st.session_state['userid'],  # user id
        "SessionId": params["session_id"],  # User session id
        "messages": [messages],  # 'messages' is a list of dictionaries
        "time": messages['time']
    }

    existing_item = DYNAMODB.Table(DYNAMODB_TABLE).get_item(Key={"UserId": st.session_state['userid'], "SessionId": params["session_id"]})

    if "Item" in existing_item:
        existing_messages = existing_item["Item"]["messages"]
        chat_item["messages"] = existing_messages + [messages]

    response = DYNAMODB.Table(DYNAMODB_TABLE).put_item(
        Item=chat_item
    )    
def save_chat_local(file_path, new_data,params):
    """Store long term chat history Local Disk"""   
    try:
        # Read the existing JSON data from the file
        with open(file_path, "r",encoding='utf-8') as file:
            existing_data = json.load(file)
        if params["session_id"] not in existing_data:
            existing_data[params["session_id"]]=[]
    except FileNotFoundError:
        # If the file doesn't exist, initialize an empty list
        existing_data = {params["session_id"]:[]}
    # Append the new data to the existing list
    from decimal import Decimal
    data = [{k: float(v) if isinstance(v, Decimal) else v for k, v in item.items()} for item in new_data]
    existing_data[params["session_id"]].extend(data)
    # Write the updated list back to the JSON file
    with open(file_path, "w", encoding="utf-8") as file:
        json.dump(existing_data, file)
        
def load_chat_local(file_path,params):
    """Load long term chat history from Local"""   
    try:
        # Read the existing JSON data from the file
        with open(file_path, "r",encoding='utf-8') as file:
            existing_data = json.load(file)
            if params["session_id"] in existing_data:
                existing_data=existing_data[params["session_id"]]
            else:
                existing_data=[]
    except FileNotFoundError:
        # If the file doesn't exist, initialize an empty list
        existing_data = []
    return existing_data


def bedrock_streemer(params,response, handler):
    text=''
    think = ""
    signature = ""
    for chunk in response['stream']:    
        if 'contentBlockDelta' in chunk:
            delta = chunk['contentBlockDelta']['delta']       
            # print(chunk)
            if 'text' in delta:
                text += delta['text']               
                handler.markdown(text.replace("$", "\\$"), unsafe_allow_html=True)
            if 'reasoningContent' in delta:
                if "text" in delta['reasoningContent']:
                    think += delta['reasoningContent']['text']                    
                    handler.markdown('**MODEL REASONING**\n\n'+ think.replace("$", "\\$"), unsafe_allow_html=True)
                elif "signature" in delta['reasoningContent']:
                    signature = delta['reasoningContent']['signature']
    
        elif "metadata" in chunk:
            
            if 'cacheReadInputTokens' in chunk['metadata']['usage']:
                print(f"\nCache Read Tokens: {chunk['metadata']['usage']['cacheReadInputTokens']}")
                print(f"Cache Write Tokens: {chunk['metadata']['usage']['cacheWriteInputTokens']}")
            input_tokens = chunk['metadata']['usage']["inputTokens"]
            output_tokens = chunk['metadata']['usage']["outputTokens"]
            latency = chunk['metadata']['metrics']["latencyMs"]
            pricing = st.session_state['input_token'] * pricing_file[f"{params['model']}"]["input"] + st.session_state['output_token'] * pricing_file[f"{params['model']}"]["output"]
            st.session_state['cost']+=pricing 
            print(f"\nInput Tokens: {input_tokens}\nOutput Tokens: {output_tokens}\nLatency: {latency}ms")
    return text, think

def bedrock_claude_(params, chat_history, system_message, prompt, model_id, image_path=None, handler=None):
    # st.write(chat_history)
    chat_history_copy = chat_history[:]
    content = []
    if image_path:
        if not isinstance(image_path, list):
            image_path = [image_path]
        for img in image_path:
            s3 = boto3.client('s3')
            match = re.match("s3://(.+?)/(.+)", img)
            image_name = os.path.basename(img)
            _, ext = os.path.splitext(image_name)
            if "jpg" in ext:
                ext = ".jpeg"
            bucket_name = match.group(1)
            key = match.group(2)
            if ".plotly" in key:
                bytes_image = plotly_to_png_bytes(img)
                ext = ".png"
            else:
                obj = s3.get_object(Bucket=bucket_name, Key=key)
                bytes_image = obj['Body'].read()
            content.extend([{"text": image_name}, {
              "image": {
                "format": f"{ext.lower().replace('.', '')}",
                "source": {"bytes": bytes_image}
              }
            }])

    content.append({       
        "text": prompt
            })
    chat_history_copy.append({"role": "user",
                              "content": content})
    system_message = [{"text": system_message}] 

    config = Config(
                read_timeout=600,  # Read timeout parameter
                retries=dict(
                    max_attempts=10  ## Handle retries
                )
            )
    bedrock_runtime = boto3.client(service_name='bedrock-runtime', region_name=REGION, config=config)

    if st.session_state['reasoning_mode'] and "3-7" in model_id :
        response = bedrock_runtime.converse_stream(messages=chat_history_copy, modelId=model_id,
                                                   inferenceConfig={"maxTokens": 18000, "temperature": 1,},
                                                   system=system_message,
                                                   additionalModelRequestFields={"thinking": {"type": "enabled", "budget_tokens": 10000}}
                                                  )
    else:
        response = bedrock_runtime.converse_stream(messages=chat_history_copy, modelId=model_id,
                                                   inferenceConfig={"maxTokens": 4000, "temperature": 0.5,},
                                                   system=system_message,
                                                  )

    answer, think = bedrock_streemer(params, response, handler) 
    return answer, think

def _invoke_bedrock_with_retries(params, current_chat, chat_template, question, model_id, image_path, handler):
    max_retries = 10
    backoff_base = 2
    max_backoff = 3  # Maximum backoff time in seconds
    retries = 0

    while True:
        try:
            response, think = bedrock_claude_(params, current_chat, chat_template, question, model_id, image_path, handler)
            return response, think
        except ClientError as e:
            if e.response['Error']['Code'] == 'ThrottlingException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
                else:
                    raise e
            elif e.response['Error']['Code'] == 'ModelStreamErrorException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
                else:
                    raise e
            elif e.response['Error']['Code'] == 'EventStreamError':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
                else:
                    raise e
            else:
                # Some other API error, rethrow
                raise

def parse_s3_uri(uri):
    """
    Parse an S3 URI and extract the bucket name and key.

    :param uri: S3 URI (e.g., 's3://bucket-name/path/to/file.txt')
    :return: Tuple of (bucket_name, key) if valid, (None, None) if invalid
    """
    pattern = r'^s3://([^/]+)/(.*)$'
    match = re.match(pattern, uri)
    if match:
        return match.groups()
    return (None, None)
    
def copy_s3_object(source_uri, dest_bucket, dest_key):
    """
    Copy an object from one S3 location to another.

    :param source_uri: S3 URI of the source object
    :param dest_bucket: Name of the destination bucket
    :param dest_key: Key to be used for the destination object
    :return: True if successful, False otherwise
    """
    s3 = boto3.client('s3')

    # Parse the source URI
    source_bucket, source_key = parse_s3_uri(source_uri)
    if not source_bucket or not source_key:
        print(f"Invalid source URI: {source_uri}")
        return False

    try:
        # Create a copy source dictionary
        copy_source = {
            'Bucket': source_bucket,
            'Key': source_key
        }

        # Copy the object
        s3.copy_object(CopySource=copy_source, Bucket=dest_bucket, Key=f"{dest_key}/{source_key}")
        return f"s3://{dest_bucket}/{dest_key}/{source_key}"

    except ClientError as e:
        print(f"An error occurred: {e}")
        raise(e)
        # return False

class LibraryInstallationDetected(Exception):
    """Exception raised when potential library installation is detected."""
    pass


def check_for_library_installs(code_string):
    # Check for pip install commands using subprocess
    if re.search(r'subprocess\.(?:check_call|run|Popen)\s*\(\s*\[.*pip.*install', code_string):
        raise LibraryInstallationDetected(f"Potential library installation detected: '{keyword}' found in code.")

    # Check for pip as a module
    if re.search(r'pip\._internal\.main\(\[.*install', code_string) or re.search(r'pip\.main\(\[.*install', code_string):
        raise LibraryInstallationDetected(f"Potential library installation detected: '{keyword}' found in code.")

    keywords = ["subprocess","pip","conda","install","easy_install","setup.py","pipenv",
            "git+","svn+","hg+","bzr+","requirements.txt","environment.yml","apt-get","yum","brew",
            "ensurepip","get-pip","pkg_resources","importlib","setuptools","distutils","venv","virtualenv",
            "pyenv"]
    
    # Convert the code string to lowercase for case-insensitive matching
    code_lower = code_string.lower()
    
    # Check for each keyword
    for keyword in keywords:
        if keyword in code_lower:
            return True

    # Check for each keyword
    for keyword in keywords:
        if keyword in code_lower:
            raise LibraryInstallationDetected(f"Potential library installation detected: '{keyword}' found in code.")
   
                
def put_obj_in_s3_bucket_(docs):
    """Uploads a file to an S3 bucket and returns the S3 URI of the uploaded object.
    Args:
       docs (str): The local file path of the file to upload to S3.
   Returns:
       str: The S3 URI of the uploaded object, in the format "s3://{bucket_name}/{file_path}".
    """
    if isinstance(docs,str):
        s3_uri_pattern = r'^s3://([^/]+)/(.*?([^/]+)/?)$'
        if bool(re.match(s3_uri_pattern,  docs)):
            file_uri=copy_s3_object(docs, BUCKET, S3_DOC_CACHE_PATH)
            return file_uri
        else:
            file_name=os.path.basename(docs)
            file_path=f"{S3_DOC_CACHE_PATH}/{docs}"
            S3.upload_file(docs, BUCKET, file_path)
            return f"s3://{BUCKET}/{file_path}"
    else:
        file_name=os.path.basename(docs.name)
        file_path=f"{S3_DOC_CACHE_PATH}/{file_name}"
        S3.put_object(Body=docs.read(),Bucket= BUCKET, Key=file_path)           
        return f"s3://{BUCKET}/{file_path}"


def get_large_s3_obj_from_bucket_(file, max_bytes=1000000):
    """Retrieves a portion of an object from an S3 bucket given its S3 URI.
    Args:
       file (str): The S3 URI of the object to retrieve, in the format "s3://{bucket_name}/{key}".
       max_bytes (int, optional): Maximum number of bytes to read from the beginning of the file.   
    Returns:
       botocore.response.StreamingBody: The retrieved S3 object.
    """
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)    
    bucket_name = match.group(1)
    key = match.group(2)

    if max_bytes:
        # Read specific number of bytes
        obj = s3.get_object(Bucket=bucket_name, Key=key, Range=f'bytes=0-{max_bytes-1}')
    else:
        # Read the whole object if max_bytes is not specified
        obj = s3.get_object(Bucket=bucket_name, Key=key)

    return obj


def get_s3_obj_from_bucket_(file):
    """Retrieves an object from an S3 bucket given its S3 URI.
    Args:
       file (str): The S3 URI of the object to retrieve, in the format "s3://{bucket_name}/{key}".
   Returns:
       botocore.response.StreamingBody: The retrieved S3 object.
    """
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)    
    bucket_name = match.group(1)
    key = match.group(2)    
    obj = s3.get_object(Bucket=bucket_name, Key=key)  
    return obj


def iter_block_items(parent):
    if isinstance(parent, Document):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError("something's not right")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield DocxTable(child, parent)

def extract_text_and_tables(docx_path):
    """ Extract text from docx files"""
    document = DocxDocument(docx_path)
    content = ""
    current_section = ""
    section_type = None
    for block in iter_block_items(document):
        if isinstance(block, Paragraph):
            if block.text:
                if block.style.name == 'Heading 1':
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                        section_type = None  
                    section_type ="h1"
                    content += f"<{section_type}>{block.text}</{section_type}>\n"
                elif block.style.name== 'Heading 3':
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                    section_type = "h3"  
                    content += f"<{section_type}>{block.text}</{section_type}>\n"
                
                elif block.style.name == 'List Paragraph':
                    # Add to the current list section
                    if section_type != "list":
                        # Close the current section if it exists
                        if current_section:
                            content += f"{current_section}</{section_type}>\n"
                        section_type = "list"
                        current_section = "<list>"
                    current_section += f"{block.text}\n"
                elif block.style.name.startswith('toc'):
                    # Add to the current toc section
                    if section_type != "toc":
                        # Close the current section if it exists
                        if current_section:
                            content += f"{current_section}</{section_type}>\n"
                        section_type = "toc"
                        current_section = "<toc>"
                    current_section += f"{block.text}\n"
                else:
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                        section_type = None
                    
                    # Append the passage text without tagging
                    content += f"{block.text}\n"
        
        elif isinstance(block, DocxTable):
            # Add the current section before the table
            if current_section:
                content += f"{current_section}</{section_type}>\n"
                current_section = ""
                section_type = None

            content += "<table>\n"
            for row in block.rows:
                row_content = []
                for cell in row.cells:
                    cell_content = []
                    for nested_block in iter_block_items(cell):
                        if isinstance(nested_block, Paragraph):
                            cell_content.append(nested_block.text)
                        elif isinstance(nested_block, DocxTable):
                            nested_table_content = parse_nested_table(nested_block)
                            cell_content.append(nested_table_content)
                    row_content.append(CSV_SEPERATOR.join(cell_content))
                content += CSV_SEPERATOR.join(row_content) + "\n"
            content += "</table>\n"

    # Add the final section
    if current_section:
        content += f"{current_section}</{section_type}>\n"

    return content

def parse_nested_table(table):
    nested_table_content = "<table>\n"
    for row in table.rows:
        row_content = []
        for cell in row.cells:
            cell_content = []
            for nested_block in iter_block_items(cell):
                if isinstance(nested_block, Paragraph):
                    cell_content.append(nested_block.text)
                elif isinstance(nested_block, DocxTable):
                    nested_table_content += parse_nested_table(nested_block)
            row_content.append(CSV_SEPERATOR.join(cell_content))
        nested_table_content += CSV_SEPERATOR.join(row_content) + "\n"
    nested_table_content += "</table>"
    return nested_table_content



def extract_text_from_pptx_s3(pptx_buffer):
    """ Extract Text from pptx files"""
    presentation = Presentation(pptx_buffer)    
    text_content = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text)
        text_content.append('\n'.join(slide_text))    
    return '\n\n'.join(text_content)

def get_s3_keys(prefix):
    """list all keys in an s3 path"""
    s3 = boto3.client('s3')
    keys = []
    next_token = None
    while True:
        if next_token:
            response = s3.list_objects_v2(Bucket=BUCKET, Prefix=prefix, ContinuationToken=next_token)
        else:
            response = s3.list_objects_v2(Bucket=BUCKET, Prefix=prefix)
        if "Contents" in response:
            for obj in response['Contents']:
                key = obj['Key']
                name = key[len(prefix):]
                keys.append(name)
        if "NextContinuationToken" in response:
            next_token = response["NextContinuationToken"]
        else:
            break
    return keys

def get_object_with_retry(bucket, key):
    max_retries=5
    retries = 0   
    backoff_base = 2
    max_backoff = 3  # Maximum backoff time in seconds
    s3 = boto3.client('s3')
    while retries < max_retries:
        try:
            response = s3.get_object(Bucket=bucket, Key=key)
            return response
        except ClientError as e:
            error_code = e.response['Error']['Code']
            if error_code == 'DecryptionFailureException':
                sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                print(f"Decryption failed, retrying in {sleep_time} seconds...")                
                time.sleep(sleep_time)               
                retries += 1
            elif e.response['Error']['Code'] == 'ModelStreamErrorException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
            else:
                raise e
                
def exract_pdf_text_aws(file):    
    file_base_name=os.path.basename(file)
    dir_name, ext = os.path.splitext(file)
    # Checking if extracted doc content is in S3
    if USE_TEXTRACT:        
        if [x for x in get_s3_keys(f"{TEXTRACT_RESULT_CACHE_PATH}/") if file_base_name in x]:    
            response = get_object_with_retry(BUCKET, f"{TEXTRACT_RESULT_CACHE_PATH}/{file_base_name}.txt")
            text = response['Body'].read().decode()
            return text
        else:            
            extractor = Textractor()
            # Asynchronous call, you will experience some wait time. Try caching results for better experience
            if "pdf" in ext:
                print("Asynchronous call, you may experience some wait time.")
                document = extractor.start_document_analysis(
                file_source=file,
                features=[TextractFeatures.LAYOUT,TextractFeatures.TABLES],       
                save_image=False,   
                s3_output_path=f"s3://{BUCKET}/textract_output/"
            )
            #Synchronous call
            else:
                document = extractor.analyze_document(
                file_source=file,
                features=[TextractFeatures.LAYOUT,TextractFeatures.TABLES],  
                save_image=False,
            )
            config = TextLinearizationConfig(
            hide_figure_layout=False,   
            hide_header_layout=False,    
            table_prefix="<table>",
            table_suffix="</table>",
            )
            # Upload extracted content to s3
            S3.put_object(Body=document.get_text(config=config), Bucket=BUCKET, Key=f"{TEXTRACT_RESULT_CACHE_PATH}/{file_base_name}.txt") 
            return document.get_text(config=config)
    else:
        s3=boto3.resource("s3")
        match = re.match("s3://(.+?)/(.+)", file)
        if match:
            bucket_name = match.group(1)
            key = match.group(2)
        if "pdf" in ext:            
            pdf_bytes = io.BytesIO()            
            s3.Bucket(bucket_name).download_fileobj(key, pdf_bytes)
            # Read the PDF from the BytesIO object
            pdf_bytes.seek(0)                      
            # Create a PDF reader object
            pdf_reader = PyPDF2.PdfReader(pdf_bytes)
            # Get the number of pages in the PDF
            num_pages = len(pdf_reader.pages)
            # Extract text from each page
            text = ''
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text += page.extract_text()
        else:
            img_bytes = io.BytesIO()
            s3.Bucket(bucket_name).download_fileobj(key, img_bytes)
            img_bytes.seek(0)         
            image_stream = io.BytesIO(image_bytes)
            image = Image.open(image_stream)
            text = pytesseract.image_to_string(image)
        return text    

def detect_encoding(s3_uri):
    """detect csv encoding"""
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", s3_uri)
    if match:
        bucket_name = match.group(1)
        key = match.group(2) 
    response = get_large_s3_obj_from_bucket_(s3_uri)
    content = response['Body'].read()
    result = chardet.detect(content)
    df = content.decode(result['encoding'])
    return result['encoding'], df

class InvalidContentError(Exception):
    pass

def parse_csv_from_s3(s3_uri):
    """Here we are only loading the first 3 rows to the model. 3 rows is sufficient for the model to figure out the schema"""
    
    try:
        # Detect the file encoding using chardet
        encoding, content = detect_encoding(s3_uri)  
        # Use StringIO to create a file-like object
        csv_file = io.StringIO(content)
        # Read the CSV file using pandas
        df = pd.read_csv(csv_file, delimiter=None, engine='python').iloc[:3]
        data_types = df.dtypes
        return f"{df.to_csv(index=False)}\n\nHere are the data types for each columns:\n{data_types}"
        
    except Exception as e:
        raise InvalidContentError(f"Error: {e}")
        
        
def strip_newline(cell):
    return str(cell).strip()

def table_parser_openpyxl(file):
    """
    Here we are only loading the first 20 rows to the model and we are not massaging the dataset by merging empty cells. 5 rows is sufficient for the model to figure out the schema
    """
    # Read from S3
    s3 = boto3.client('s3', region_name=REGION)
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
        obj = s3.get_object(Bucket=bucket_name, Key=key)    
        # Read Excel file from S3 into a buffer
        xlsx_buffer = io.BytesIO(obj['Body'].read()) 
        # Load workbook
        wb=pd.read_excel(xlsx_buffer,sheet_name=None, header=None)  
        all_sheets_string=""
        # Iterate over each sheet in the workbook
        for sheet_name, sheet_data in wb.items():
            df = pd.DataFrame(sheet_data)
            # Convert to string and tag by sheet name
            all_sheets_string+=f'Here is a data preview of this sheet (first 5 rows):<{sheet_name}>\n{df.iloc[:5].to_csv(index=False, header=False)}\n</{sheet_name}>\n'
        return all_sheets_string
    else:
        raise Exception(f"{file} not formatted as an S3 path")

def calamaine_excel_engine(file):
    """
    Here we are only loading the first 20 rows to the model and we are not massaging the dataset by merging empty cells. 20 rows is sufficient for the model to figure out the schema
    """
    # # Read from S3
    s3 = boto3.client('s3',region_name=REGION)
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
        obj = s3.get_object(Bucket=bucket_name, Key=key)
        # Read Excel file from S3 into a buffer
        xlsx_buffer = io.BytesIO(obj['Body'].read())
        xlsx_buffer.seek(0)
        all_sheets_string = ""
        # Load the Excel file
        workbook = CalamineWorkbook.from_filelike(xlsx_buffer)
        # Iterate over each sheet in the workbook
        for sheet_name in workbook.sheet_names:
            # Get the sheet by name
            sheet = workbook.get_sheet_by_name(sheet_name)
            df = pd.DataFrame(sheet.to_python(skip_empty_area=False))
            df = df.map(strip_newline)
            all_sheets_string += f'Here is a data preview of this sheet (first 5 rows):\n\n<{sheet_name}>\n{df.iloc[:5].to_csv(index=False, header=0)}\n</{sheet_name}>\n'
        return all_sheets_string
    else:
        raise Exception(f"{file} not formatted as an S3 path")

def table_parser_utills(file):
    try:
        response = table_parser_openpyxl(file)
        if response:
            return response
        else:
            return calamaine_excel_engine(file)
    except Exception:
        try:
            return calamaine_excel_engine(file)
        except Exception as e:
            raise Exception(str(e))

def process_document_types(file):
    """Handle various document format"""
    dir_name, ext = os.path.splitext(file)
    if ".csv" == ext.lower():
        content = parse_csv_from_s3(file)
    elif ext.lower() in [".txt", ".py"]:
        obj = get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
    elif ext.lower() in [".xlsx", ".xls"]:
        content = table_parser_utills(file)
    elif ext.lower() in [".pdf", ".png", ".jpg", ".tif", ".jpeg"]:
        content = exract_pdf_text_aws(file)
    elif ".json" == ext.lower():
        obj = get_s3_obj_from_bucket_(file)
        content = json.loads(obj['Body'].read())
    elif ".docx" == ext.lower():
        obj = get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        docx_buffer = io.BytesIO(content)
        content = extract_text_and_tables(docx_buffer)
    elif ".pptx" == ext.lower():
        obj = get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        docx_buffer = io.BytesIO(content)
        content = extract_text_from_pptx_s3(docx_buffer)

    # Implement any other file extension logic
    return content

def plotly_to_png_bytes(s3_uri):
    """
    Read a .plotly file from S3 given an S3 URI, convert it to a PNG image, and return the image as bytes.

    :param s3_uri: S3 URI of the .plotly file (e.g., 's3://bucket-name/path/to/file.plotly')
    :return: PNG image as bytes
    """
    # Parse S3 URI
    parsed_uri = urlparse(s3_uri)
    bucket_name = parsed_uri.netloc
    file_key = parsed_uri.path.lstrip('/')

    # Initialize S3 client
    s3_client = boto3.client('s3')

    try:
        # Read the .plotly file from S3
        response = s3_client.get_object(Bucket=bucket_name, Key=file_key)
        plotly_data = json.loads(response['Body'].read().decode('utf-8'))

        # Create a Figure object from the plotly data
        fig = go.Figure(data=plotly_data['data'], layout=plotly_data.get('layout', {}))

        # Convert the figure to PNG bytes
        img_bytes = fig.to_image(format="png")

        return img_bytes

    except Exception as e:
        print(f"An error occurred: {str(e)}")
        return None


def get_chat_history_db(params, cutoff,vision_model):
    """
    Load chat history and attachments from DynamoDB and S3 accordingly

    parameters:
    params (dict): Application parameters
    cutoff (int): Custoff of Chat history to be loaded
    vision_model (bool): Boolean if Claude 3 model is used 
    """
    current_chat, chat_hist = [], []
    if params['chat_histories'] and cutoff != 0:
        chat_hist = params['chat_histories'][-cutoff:]
        for ids, d in enumerate(chat_hist):
            if d['image'] and vision_model and LOAD_DOC_IN_ALL_CHAT_CONVO:
                content = []
                for img in d['image']:
                    s3 = boto3.client('s3')
                    match = re.match("s3://(.+?)/(.+)", img)
                    image_name = os.path.basename(img)
                    _, ext = os.path.splitext(image_name)
                    if "jpg" in ext:
                        ext = ".jpeg"
                    if match:
                        bucket_name = match.group(1)
                        key = match.group(2)
                        if ".plotly" in key:
                            bytes_image = plotly_to_png_bytes(img)
                            ext = ".png"
                        else:
                            obj = s3.get_object(Bucket=bucket_name, Key=key)
                            bytes_image = obj['Body'].read()
                    content.extend(
                        [
                            {"text": image_name},
                            {'image':
                                {
                                    'format': ext.lower().replace('.', ''),
                                    'source': {'bytes': bytes_image}
                                        }
                             }
                        ]
                    )
                content.extend([{"text": d['user']}])
                if 'tool_result_id' in d and d['tool_result_id']:
                    user = [{'toolResult': {'toolUseId': d['tool_result_id'],
                             'content': content}}]
                    current_chat.append({'role': 'user', 'content': user})
                else:
                    current_chat.append({'role': 'user', 'content': content})
            elif d['document'] and LOAD_DOC_IN_ALL_CHAT_CONVO: 
                doc = 'Here is a document showing sample rows:\n'
                for docs in d['document']:
                    uploads = process_document_types(docs)
                    doc_name = os.path.basename(docs)
                    doc += f"<{doc_name}>\n{uploads}\n</{doc_name}>\n"
                if not vision_model and d["image"]:
                    for docs in d['image']:
                        uploads = process_document_types(docs)
                        doc_name = os.path.basename(docs)
                        doc += f"<{doc_name}>\n{uploads}\n</{doc_name}>\n"
                current_chat.append({'role': 'user', 'content': [{"text": doc + d['user']}]})
                # do not have a tool return for the document section because the tool does not provide documents only images
            else:
                if 'tool_result_id' in d and d['tool_result_id']:
                    user = [{'toolResult': {'toolUseId': d['tool_result_id'],
                             'content': [{'text': d['user']}]}}]
                    current_chat.append({'role': 'user', 'content': user})
                else:
                    current_chat.append({'role': 'user', 'content': [{"text": d['user']}]})

            if 'tool_use_id' in d and d['tool_use_id']:
                assistant = [{'toolUse': {'toolUseId': d['tool_use_id'],
                                          'name': d['tool_name'],
                                          'input': {'code': d['assistant'],
                                                    "dataset_name": d['tool_params']['ds'],
                                                    "python_packages": d['tool_params']['pp']}}}
                             ]

                current_chat.append({'role': 'assistant', 'content': assistant})
            else:
                current_chat.append({'role': 'assistant', 'content': [{"text": d['assistant']}]})
    return current_chat, chat_hist

def stream_messages(bedrock_client,
                    model_id,
                    messages,
                    tool_config,
                    system,
                    temperature,
                    handler):
    """
    Sends a message to a model and streams the response.
    Args:
        bedrock_client: The Boto3 Bedrock runtime client.
        model_id (str): The model ID to use.
        messages (JSON) : The messages to send to the model.
        tool_config : Tool Information to send to the model.

    Returns:
        stop_reason (str): The reason why the model stopped generating text.
        message (JSON): The message that the model generated.

    """

    if st.session_state['reasoning_mode']:
        response = bedrock_client.converse_stream(
            modelId=model_id,
            messages=messages,
            inferenceConfig={"maxTokens": 18000, "temperature": 1},
            toolConfig=tool_config,
            system=system,
            additionalModelRequestFields={"thinking": {"type": "enabled", "budget_tokens": 10000}}
        )
    else:
        response = bedrock_client.converse_stream(
                modelId=model_id,
                messages=messages,
                inferenceConfig={"maxTokens": 8000, "temperature": temperature},
                toolConfig=tool_config,
                system=system
        )

    stop_reason = ""
    message = {}
    content = []
    message['content'] = content
    text = ''
    tool_use = {}
    think = ''
    signature = ''

    for chunk in response['stream']:
        if 'messageStart' in chunk:
            message['role'] = chunk['messageStart']['role']
        elif 'contentBlockStart' in chunk:
            tool = chunk['contentBlockStart']['start']['toolUse']
            tool_use['toolUseId'] = tool['toolUseId']
            tool_use['name'] = tool['name']
        elif 'contentBlockDelta' in chunk:
            delta = chunk['contentBlockDelta']['delta']
            if 'toolUse' in delta:
                if 'input' not in tool_use:
                    tool_use['input'] = ''
                tool_use['input'] += delta['toolUse']['input']
            elif 'text' in delta:
                text += delta['text']
                if handler:
                    handler.markdown(text.replace("$", "\\$"), unsafe_allow_html=True)
            elif 'reasoningContent' in delta:
                if "text" in delta['reasoningContent']:
                    think += delta['reasoningContent']['text']
                    handler.markdown('**MODEL REASONING**\n\n' + think.replace("$", "\\$"), unsafe_allow_html=True)
                if "signature" in delta['reasoningContent']:
                    signature = delta['reasoningContent']['signature']

        elif 'contentBlockStop' in chunk:
            if 'input' in tool_use:
                tool_use['input'] = json.loads(tool_use['input'])
                content.append({'toolUse': tool_use})
            else:
                content.append({'text': text})
                text = ''

        elif 'messageStop' in chunk:
            stop_reason = chunk['messageStop']['stopReason']
        elif "metadata" in chunk:
            input_tokens = chunk['metadata']['usage']["inputTokens"]
            output_tokens = chunk['metadata']['usage']["outputTokens"]
            latency = chunk['metadata']['metrics']["latencyMs"]

    if tool_use:
        try:
            handler.markdown(f"{text}\n```python\n{message['content'][1]['toolUse']['input']['code']}", unsafe_allow_html=True )
        except:
            if len(message['content']) == 3:
                handler.markdown(f"{text}\n```python\n{message['content'][2]['toolUse']['input']['code']}", unsafe_allow_html=True )
            else:
                handler.markdown(f"{text}\n```python\n{message['content'][0]['toolUse']['input']['code']}", unsafe_allow_html=True )
    return stop_reason, message, input_tokens, output_tokens, think


def self_crtique(params, code, error, dataset, handler=None):

    import re
    if params["engine"] == "pyspark":
        with open("prompt/pyspark_debug_prompt.txt", "r") as fo:
            prompt_template = fo.read()
        values = {"dataset": dataset, "code": code, "error": error}
        prompt = prompt_template.format(**values)
        system = "You are an expert pyspark debugger for Amazon Athena PySpark Runtime"
    else:
        with open("prompt/python_debug_prompt.txt", "r") as fo:
            prompt_template = fo.read()
        values = {"dataset": dataset, "code": code, "error": error}
        prompt = prompt_template.format(**values)
        system = "You are an expert python debugger."

    model_id = 'us.' + model_info[params['model']]
    fixed_code, think = _invoke_bedrock_with_retries(params, [], system, prompt, model_id, [], handler)
    code_pattern = r'<code>(.*?)</code>'
    match = re.search(code_pattern, fixed_code, re.DOTALL)
    code = match.group(1)
    if handler:
        handler.markdown(f"```python\n{code}", unsafe_allow_html=True)
    lib_pattern = r'<package>(.*?)</package>'
    match = re.search(lib_pattern, fixed_code, re.DOTALL)
    if match:
        libs = match.group(1)
    else:
        libs = ''
    return code, libs

def process_files(files):
    result_string = ""
    errors = []
    future_proxy_mapping = {}
    futures = []

    with concurrent.futures.ProcessPoolExecutor() as executor:
        # Partial function to pass the process_document_types function
        func = partial(process_document_types)
        for file in files:
            future = executor.submit(func, file)
            future_proxy_mapping[future] = file
            futures.append(future)

        # Collect the results and handle exceptions
        for future in concurrent.futures.as_completed(futures):
            file_url = future_proxy_mapping[future]
            try:
                result = future.result()
                doc_name = os.path.basename(file_url)
                result_string += f"<s3://{BUCKET}/{S3_DOC_CACHE_PATH}/{doc_name}>\n{result}\n</s3://{BUCKET}/{S3_DOC_CACHE_PATH}/{doc_name}>\n"
            except Exception as e:
                # Get the original function arguments from the Future object
                error = {'file': file_url, 'error': str(e)}
                errors.append(error)

    return errors, result_string

def invoke_lambda(function_name, payload):
    config = Config(
            connect_timeout=600,
            read_timeout=600,  # Read timeout parameter
            retries=dict(
                max_attempts=0,  # Handle retries
                total_max_attempts=1
            )
        )
    lambda_client = boto3.client('lambda', config=config)
    response = lambda_client.invoke(
        FunctionName=function_name,
        InvocationType='RequestResponse',
        Payload=json.dumps(payload)
    )
    return json.loads(response['Payload'].read().decode('utf-8'))

class CodeExecutionError(Exception):
    pass

def load_json_data(output_content):
    try:
        json_data = json.loads(output_content.replace("'", '"'))
        return json_data
    except json.JSONDecodeError:
        try:
            parsed_data = ast.literal_eval(output_content)
            json_data = json.loads(json.dumps(parsed_data))
            return json_data
        except (SyntaxError, ValueError):
            try:
                # Replace escaped single quotes with double quotes, but handle nested quotes carefully
                modified_content = output_content
                # First, ensure the outer structure is properly quoted
                if modified_content.startswith("'") and modified_content.endswith("'"):
                    modified_content = modified_content[1:-1]  # Remove outer quotes
                # Replace escaped single quotes with double quotes
                modified_content = modified_content.replace("\'", '"')
                json_data = json.loads(modified_content)
                return json_data
            except json.JSONDecodeError as e:
                print(f"Error parsing JSON using all methods: {e}")
                raise e

def function_caller_claude_(params, handler=None):
    """
    Entrypoint for streaming tool use example.
    """
    current_chat, chat_hist = get_chat_history_db(params, CHAT_HISTORY_LENGTH, True)
    if current_chat and 'toolResult' in current_chat[0]['content'][0]:
        if 'toolUseId' in current_chat[0]['content'][0]['toolResult']:
            del current_chat[0:2]

    vision_model = True
    model_id = 'us.' + model_info[params['model']]
    if any(keyword in params['model'] for keyword in TEXT_ONLY_MODELS):
        vision_model = False

    full_doc_path = []
    image_path = []
    for ids, docs in enumerate(params['upload_doc']):
        file_name = docs.name
        _, extensions = os.path.splitext(file_name)
        s3_file_name = put_obj_in_s3_bucket_(docs)
        if extensions.lower() in [".jpg", ".jpeg", ".png", ".gif", ".webp"] and vision_model:
            image_path.append(s3_file_name)
            continue
        full_doc_path.append(s3_file_name)

    if params['s3_objects']:
        for ids, docs in enumerate(params['s3_objects']):
            file_name = docs
            _, extensions = os.path.splitext(file_name)
            docs = f"s3://{INPUT_BUCKET}/{INPUT_S3_PATH}/{docs}"
            full_doc_path.append(docs)
            if extensions.lower() in [".jpg", ".jpeg", ".png", ".gif", ".webp"] and vision_model:
                image_path.append(docs)
                continue
  
    errors, result_string = process_files(full_doc_path) 
    if errors:
        st.error(errors)
    question = params['question']
    if result_string and ('.csv' in result_string or '.parquet' in result_string or '.xlsx' in result_string):
        input_text = f"Here is a subset (first few rows) of the data from each dataset tagged by each file name:\n{result_string}\n{question}"
    elif result_string and not ('.csv' in result_string or '.xlsx' in result_string or '.parquet' in result_string):
        doc = 'I have provided documents and/or images tagged by their file names:\n'
        input_text = f"{doc}{result_string}\n{question}"
    else:
        input_text = question
    bedrock_client = boto3.client(service_name='bedrock-runtime', region_name=REGION, config=config)
    # Create the initial message from the user input.
    content = []
    if image_path:
        for img in image_path:
            s3 = boto3.client('s3')
            match = re.match("s3://(.+?)/(.+)", img)
            image_name = os.path.basename(img)
            _, ext = os.path.splitext(image_name)
            if "jpg" in ext:
                ext = ".jpeg"
            bucket_name = match.group(1)
            key = match.group(2)
            if ".plotly" in key:
                bytes_image = plotly_to_png_bytes(img)
                ext = ".png"
            else:
                obj = s3.get_object(Bucket=bucket_name, Key=key)
                bytes_image = obj['Body'].read()
            content.extend([{"text": image_name}, {
              "image": {
                "format": f"{ext.lower().replace('.', '')}",
                "source": {"bytes": bytes_image}
              }
            }])
    content.append({"text": input_text})
    messages = [{
        "role": "user",
        "content": content
    }]
    # Define the tool and prompt template to send to the model.
    if params["engine"] == "pyspark":
        with open("prompt/pyspark_tool_template.json", "r") as f:
            tool_config = json.load(f)
        with open("prompt/pyspark_tool_prompt.txt", "r") as fo:
            description = fo.read()
        with open("prompt/pyspark_tool_system.txt", "r") as fod:
            system_prompt = fod.read()
        tool_config['tools'][0]['toolSpec']['inputSchema']['json']['properties']['code']['description'] = description
    else:
        with open("prompt/python_tool_template.json", "r") as f:
            tool_config = json.load(f)
        with open("prompt/python_tool_prompt.txt", "r") as fo:
            description = fo.read()
        with open("prompt/python_tool_system.txt", "r") as fod:
            system_prompt = fod.read()
        tool_config['tools'][0]['toolSpec']['inputSchema']['json']['properties']['code']['description'] = description

    system = [
            {
                'text': system_prompt
            }
        ]
    current_chat.extend(messages)
    # Send the message and get the tool use request from response.
    stop_reason, message, input_tokens, output_tokens, think = stream_messages(
        bedrock_client, model_id, current_chat, tool_config, system, 0.1, handler)
    messages.append(message)
    if stop_reason != "tool_use":
        chat_history = {
                        "user": question,
                        "assistant": message['content'][0]['text'] if message['content'][0]['text'] else message['content'][1]['text'],
                        "image": image_path,
                        "document": full_doc_path,
                        "thinking": think,
                        "modelID": model_id,
                        "time": str(time.time()),
                        "input_token": round(input_tokens),
                        "output_token": round(output_tokens)
                     }
        if DYNAMODB_TABLE:
            put_db(params, chat_history)
        # use local disk for storage
        else:
            save_chat_local(LOCAL_CHAT_FILE_NAME, [chat_history], params)

        return message['content'][0]['text'], "", "", "", full_doc_path, stop_reason, ""
    elif stop_reason == "tool_use":

        self_correction_retry = 5
        for content in message['content']:
            if 'toolUse' in content:
                tool = content['toolUse']
                if tool['name'] == tool_config['tools'][0]['toolSpec']['name']:
                    # Preliminary Guardrail to check that the code does not have any install commands 
                    check_for_library_installs(tool['input']['code'])
                    i = 0
                    while i < self_correction_retry:
                        try:
                            payload = {
                                    "python_packages": tool['input']['python_packages'],
                                    "code": tool['input']['code'],
                                    "dataset_name": tool['input']['dataset_name'],
                                    "iterate": i,
                                    "bucket": BUCKET,
                                    "file_path": S3_DOC_CACHE_PATH
                                }
                            if params["engine"] == "pyspark":
                                tool_execution_response = send_athena_job(payload, ATHENA_WORKGROUP_NAME)  # Execute generated code in Amazon Athena
                                if "error" not in tool_execution_response:
                                    pattern = r'<output>(.*?)</output>'
                                    match = re.search(pattern, tool_execution_response['stdout'], re.DOTALL)
                                    if match:
                                        output_content = match.group(1)
                                        # Parse the extracted content to JSON
                                        try:
                                            json_data = load_json_data(output_content)
                                            image_holder = json_data.get('plotly', [])
                                            results = json_data['result']
                                            plotly_obj = json_data.get('plotly', [])
                                        except json.JSONDecodeError as e:
                                            print(f"Error parsing JSON: {e}")
                                            raise e
                                    else:
                                        print("No <output> tags found in the input string.")
                                    break
                                else:
                                    raise Exception(tool_execution_response)
                            else:
                                # payload['code'] = tool['input']['code']
                                tool_execution_response = invoke_lambda(LAMBDA_FUNC, payload)  # Execute generated code in AWS Lambda
                                if tool_execution_response.get('statusCode') == 200:
                                    json_data = json.loads(tool_execution_response['body'])
                                    image_holder = json_data.get('image_dict', [])
                                    results = json_data['result']
                                    plotly_obj = json_data.get('plotly', [])
                                    break
                                else:
                                    raise Exception(tool_execution_response.get('body'))
                        except Exception as err:
                            print(f"ERROR: {err}")
                            with st.spinner(f'**Self Correction {i+1}**'):
                                tool['input']['code'], tool['input']['python_packages'] = self_crtique(params, tool['input']['code'], err, result_string, handler)
                            i += 1
                    if i == self_correction_retry:
                        raise CodeExecutionError("Request Failed due to exceed on self-correction trials")

        if 'text' in message['content'][0] and message['content'][0]['text']:
            code = tool['input']['code']
            ds = message['content'][1]['toolUse']['input']['dataset_name']
            pp = message['content'][1]['toolUse']['input']['python_packages']
            tool_ids = message['content'][1]['toolUse']['toolUseId']
            tool_name = message['content'][1]['toolUse']['name']
        else:
            try:
                code = tool['input']['code']
                ds = message['content'][0]['toolUse']['input']['dataset_name']
                pp = message['content'][0]['toolUse']['input']['python_packages']
                tool_ids = message['content'][0]['toolUse']['toolUseId']
                tool_name = message['content'][0]['toolUse']['name']
            except Exception:
                code = tool['input']['code']
                ds = message['content'][-1]['toolUse']['input']['dataset_name']
                pp = message['content'][-1]['toolUse']['input']['python_packages']
                tool_ids = message['content'][-1]['toolUse']['toolUseId']
                tool_name = message['content'][-1]['toolUse']['name']

        chat_history = {"user": question,
                        "assistant": code,
                        "image": [],
                        "document": full_doc_path,
                        "modelID": model_id,
                        "time": str(time.time()),
                        "input_token": round(input_tokens),
                        "output_token": round(output_tokens),
                        "tool_use_id": tool_ids,
                        "tool_name": tool_name,
                        "tool_params": {"ds": ds, "pp": pp}
                        }

        if DYNAMODB_TABLE:
            put_db(params, chat_history)
        # use local disk for storage
        else:
            save_chat_local(LOCAL_CHAT_FILE_NAME, [chat_history], params)
        return "", tool, results, image_holder, full_doc_path, stop_reason, plotly_obj